from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import io

# ── Design Palette ────────────────────────────────────────────────────────────
DARK_BG      = RGBColor(0x0D, 0x1B, 0x2A)   # Deep navy
ACCENT_BLUE  = RGBColor(0x00, 0x7B, 0xFF)   # Vivid blue
ACCENT_CYAN  = RGBColor(0x00, 0xD4, 0xFF)   # Cyan highlight
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xD6, 0xE8)
MID_GRAY     = RGBColor(0x8A, 0x9B, 0xB8)
CARD_BG      = RGBColor(0x13, 0x2A, 0x45)   # Slightly lighter navy

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _set_bg(slide, color: RGBColor):
    """Fill slide background with a solid colour."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color: RGBColor, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_text(slide, text, left, top, width, height,
              font_size=18, bold=False, color=WHITE,
              align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def _add_bullet_slide(prs, title_text, bullets, subtitle=None):
    """Generic dark slide with a title, optional subtitle, and bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide, DARK_BG)

    # Left accent bar
    _add_rect(slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_BLUE)

    # Top gradient bar
    _add_rect(slide, Inches(0.08), Inches(0), SLIDE_W - Inches(0.08), Inches(0.06), ACCENT_CYAN)

    # Title
    _add_text(slide, title_text,
              Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.8),
              font_size=32, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.LEFT)

    # Divider line (thin rect)
    _add_rect(slide, Inches(0.4), Inches(1.0), Inches(11.5), Inches(0.03), ACCENT_BLUE)

    if subtitle:
        _add_text(slide, subtitle,
                  Inches(0.4), Inches(1.05), Inches(12.5), Inches(0.5),
                  font_size=14, color=MID_GRAY, italic=True)
        bullet_top = Inches(1.6)
    else:
        bullet_top = Inches(1.2)

    # Bullet points
    y = bullet_top
    for i, bullet in enumerate(bullets):
        if not bullet:
            continue
        # Bullet card background
        _add_rect(slide, Inches(0.4), y, Inches(12.4), Inches(0.52), CARD_BG)
        # Bullet dot
        _add_rect(slide, Inches(0.55), y + Inches(0.18), Inches(0.12), Inches(0.12), ACCENT_CYAN)
        # Bullet text
        _add_text(slide, str(bullet),
                  Inches(0.85), y + Inches(0.05), Inches(11.8), Inches(0.45),
                  font_size=14, color=WHITE)
        y += Inches(0.6)
        if y > Inches(7.1):
            break  # prevent overflow

    return slide


def generate_ppt(project_data):
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    title_text  = project_data.get('title', 'Project Presentation')
    domain      = project_data.get('domain', 'Engineering')
    abstract    = project_data.get('abstract', '')
    problem     = project_data.get('problem_statement', '')
    features    = project_data.get('features', [])
    arch        = project_data.get('architecture_description', '')
    tech        = project_data.get('tech_stack_details', {})
    lit_survey  = project_data.get('literature_survey', '')
    methodology = project_data.get('methodology', '')
    viva_qs     = project_data.get('viva_questions', [])

    # ── SLIDE 1 : Title Slide ─────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide, DARK_BG)

    # Big diagonal accent shape
    _add_rect(slide, Inches(9), Inches(0), Inches(4.33), SLIDE_H, CARD_BG)
    _add_rect(slide, Inches(0), Inches(5.8), SLIDE_W, Inches(0.06), ACCENT_CYAN)

    # Branding dot
    _add_rect(slide, Inches(0.5), Inches(0.5), Inches(0.5), Inches(0.5), ACCENT_BLUE)
    _add_text(slide, "AXIONX", Inches(1.1), Inches(0.45), Inches(4), Inches(0.5),
              font_size=14, bold=True, color=ACCENT_CYAN)

    # Main title
    _add_text(slide, title_text,
              Inches(0.5), Inches(1.8), Inches(8.2), Inches(2.5),
              font_size=40, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Domain tag
    _add_rect(slide, Inches(0.5), Inches(4.5), Inches(2.5), Inches(0.45), ACCENT_BLUE)
    _add_text(slide, domain.upper(),
              Inches(0.55), Inches(4.52), Inches(2.4), Inches(0.4),
              font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    _add_text(slide, "Final Year Project Presentation",
              Inches(0.5), Inches(5.1), Inches(8), Inches(0.5),
              font_size=16, color=MID_GRAY)

    _add_text(slide, "Powered by AxionX AI",
              Inches(0.5), Inches(6.8), Inches(8), Inches(0.4),
              font_size=11, color=MID_GRAY, italic=True)

    # ── SLIDE 2 : Abstract ────────────────────────────────────────────────────
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide2, DARK_BG)
    _add_rect(slide2, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_BLUE)
    _add_rect(slide2, Inches(0.08), Inches(0), SLIDE_W - Inches(0.08), Inches(0.06), ACCENT_CYAN)
    _add_text(slide2, "Abstract",
              Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.8),
              font_size=32, bold=True, color=ACCENT_CYAN)
    _add_rect(slide2, Inches(0.4), Inches(1.0), Inches(11.5), Inches(0.03), ACCENT_BLUE)
    # Abstract text box
    _add_rect(slide2, Inches(0.4), Inches(1.1), Inches(12.4), Inches(5.8), CARD_BG)
    _add_text(slide2, abstract[:1200],
              Inches(0.6), Inches(1.2), Inches(12.0), Inches(5.6),
              font_size=15, color=LIGHT_GRAY)

    # ── SLIDE 3 : Problem Statement ───────────────────────────────────────────
    if problem:
        slide3 = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(slide3, DARK_BG)
        _add_rect(slide3, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_BLUE)
        _add_rect(slide3, Inches(0.08), Inches(0), SLIDE_W - Inches(0.08), Inches(0.06), ACCENT_CYAN)
        _add_text(slide3, "Problem Statement",
                  Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.8),
                  font_size=32, bold=True, color=ACCENT_CYAN)
        _add_rect(slide3, Inches(0.4), Inches(1.0), Inches(11.5), Inches(0.03), ACCENT_BLUE)
        _add_rect(slide3, Inches(0.4), Inches(1.1), Inches(12.4), Inches(5.8), CARD_BG)
        _add_text(slide3, problem[:1200],
                  Inches(0.6), Inches(1.2), Inches(12.0), Inches(5.6),
                  font_size=15, color=LIGHT_GRAY)

    # ── SLIDE 4 : Key Features ────────────────────────────────────────────────
    if features:
        feature_bullets = [str(f) for f in features[:8]]
        _add_bullet_slide(prs, "Key Features", feature_bullets,
                          subtitle="What makes this project stand out")

    # ── SLIDE 5 : System Architecture ────────────────────────────────────────
    if arch:
        arch_slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(arch_slide, DARK_BG)
        _add_rect(arch_slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_BLUE)
        _add_rect(arch_slide, Inches(0.08), Inches(0), SLIDE_W - Inches(0.08), Inches(0.06), ACCENT_CYAN)
        _add_text(arch_slide, "System Architecture",
                  Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.8),
                  font_size=32, bold=True, color=ACCENT_CYAN)
        _add_rect(arch_slide, Inches(0.4), Inches(1.0), Inches(11.5), Inches(0.03), ACCENT_BLUE)
        _add_rect(arch_slide, Inches(0.4), Inches(1.1), Inches(12.4), Inches(5.8), CARD_BG)
        _add_text(arch_slide, arch[:1400],
                  Inches(0.6), Inches(1.2), Inches(12.0), Inches(5.6),
                  font_size=14, color=LIGHT_GRAY)

    # ── SLIDE 6 : Technology Stack ────────────────────────────────────────────
    if tech:
        tech_bullets = [f"{k.capitalize()}: {v}" for k, v in tech.items() if v]
        _add_bullet_slide(prs, "Technology Stack", tech_bullets,
                          subtitle="Tools and frameworks powering this system")

    # ── SLIDE 7 : Literature Survey ───────────────────────────────────────────
    if lit_survey:
        lit_slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_bg(lit_slide, DARK_BG)
        _add_rect(lit_slide, Inches(0), Inches(0), Inches(0.08), SLIDE_H, ACCENT_BLUE)
        _add_rect(lit_slide, Inches(0.08), Inches(0), SLIDE_W - Inches(0.08), Inches(0.06), ACCENT_CYAN)
        _add_text(lit_slide, "Literature Survey",
                  Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.8),
                  font_size=32, bold=True, color=ACCENT_CYAN)
        _add_rect(lit_slide, Inches(0.4), Inches(1.0), Inches(11.5), Inches(0.03), ACCENT_BLUE)
        _add_rect(lit_slide, Inches(0.4), Inches(1.1), Inches(12.4), Inches(5.8), CARD_BG)
        _add_text(lit_slide, lit_survey[:1200],
                  Inches(0.6), Inches(1.2), Inches(12.0), Inches(5.6),
                  font_size=14, color=LIGHT_GRAY)

    # ── SLIDE 8 : Methodology ─────────────────────────────────────────────────
    if methodology:
        steps = [s.strip() for s in methodology.split('\n') if s.strip()][:7]
        if not steps:
            steps = [methodology[:800]]
        _add_bullet_slide(prs, "Methodology", steps,
                          subtitle="Development approach and process")

    # ── SLIDE 9 : Sample Viva Questions ──────────────────────────────────────
    if viva_qs:
        q_bullets = [f"Q: {q.get('question', '')}" for q in viva_qs[:5]]
        _add_bullet_slide(prs, "Anticipated Viva Questions", q_bullets,
                          subtitle="Key questions and topics to prepare")

    # ── SLIDE 10 : Conclusion ─────────────────────────────────────────────────
    slide_end = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide_end, DARK_BG)
    _add_rect(slide_end, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_CYAN)
    _add_rect(slide_end, Inches(0), SLIDE_H - Inches(0.06), SLIDE_W, Inches(0.06), ACCENT_BLUE)
    _add_rect(slide_end, Inches(4.5), Inches(2.0), Inches(4.33), Inches(0.06), ACCENT_BLUE)

    _add_text(slide_end, "Conclusion",
              Inches(0), Inches(1.0), SLIDE_W, Inches(1.0),
              font_size=40, bold=True, color=ACCENT_CYAN, align=PP_ALIGN.CENTER)
    _add_text(slide_end, "This project presents a comprehensive solution with",
              Inches(1), Inches(2.3), Inches(11.33), Inches(0.5),
              font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    _add_text(slide_end, "professional architecture and industry-ready implementation.",
              Inches(1), Inches(2.85), Inches(11.33), Inches(0.5),
              font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    _add_text(slide_end, "Thank You",
              Inches(0), Inches(4.0), SLIDE_W, Inches(1.5),
              font_size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _add_text(slide_end, "Questions & Discussion",
              Inches(0), Inches(5.5), SLIDE_W, Inches(0.7),
              font_size=20, color=MID_GRAY, align=PP_ALIGN.CENTER)

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer
